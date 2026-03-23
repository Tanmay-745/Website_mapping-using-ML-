from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title_text
        subtitle.text = subtitle_text

    def add_content_slide(title_text, bullet_points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        tf = slide.placeholders[1].text_frame
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    # Slide 1: Title
    add_title_slide("AI Legal Portal", "A Unified Intelligence Gateway for Legal Automation\nCodebase Overview & System Architecture")

    # Slide 2: Executive Summary
    add_content_slide("Executive Summary", [
        "The AI Legal Portal is a comprehensive suite of tools designed to streamline legal operations.",
        "Unified entry point (Gateway) for multiple specialized applications.",
        "Focus on automation: notice generation, data mapping, and document tracking.",
        "Modern tech stack ensuring scalability and rapid deployment."
    ])

    # Slide 3: System Architecture
    add_content_slide("System Architecture", [
        "Monorepo Structure: Consolidates frontend and backend components.",
        "Unified Gateway (index.html): Routes users to specific sub-apps.",
        "Routing Engine (vercel.json): Manages paths for Dashboard, Mapping, Barcode, and Cover Letter tools.",
        "Backend (FastAPI): High-performance Python API for core logic."
    ])

    # Slide 4: Central Dashboard
    add_content_slide("Central Dashboard", [
        "Role: Command center for legal activities.",
        "Features: Real-time monitoring, advocate management, and status tracking.",
        "Built with React & Vite for a responsive user experience.",
        "Integrated with the unified backend for data consistency."
    ])

    # Slide 5: Legal Mapping (ML Powered)
    add_content_slide("Legal Mapping & ML", [
        "Core Feature: Intelligent field detection for CSV/Excel uploads.",
        "Uses Machine Learning (Python/Jupyter) for semantic mapping.",
        "Automates the extraction of key legal entities from raw data.",
        "Reduces manual data entry errors and processing time."
    ])

    # Slide 6: Specialized Tools
    add_content_slide("Barcode & Cover Letter App", [
        "Barcode Generator: Generates unique tracking IDs for physical documents.",
        "Cover Letter App: Bulk generation of merged PDFs for legal notices.",
        "Streamlines the transition from digital data to physical mailings.",
        "Supports high-volume document processing workflows."
    ])

    # Slide 7: Tech Stack & Deployment
    add_content_slide("Tech Stack & Deployment", [
        "Frontend: HTML5, CSS3 (Modern UI), React, Vite.",
        "Backend: Python, FastAPI.",
        "Intelligence: Machine Learning (Semantic Mapping).",
        "Cloud: Vercel for serverless hosting and routing.",
        "Version Control: Git-based development workflow."
    ])

    # Slide 8: Future Roadmap
    add_content_slide("Future Roadmap", [
        "Enhanced ML models for broader legal document classification.",
        "Mobile-responsive dashboard improvements.",
        "Native integration with court digital filing systems.",
        "Expanded bulk automation for complex legal filings."
    ])

    # Slide 9: Conclusion
    add_title_slide("Thank You", "Questions? \nAI Legal Portal: Streamlining the Future of Law")

    # Save
    prs.save('AI_Legal_Portal_Presentation.pptx')
    print("Presentation created successfully as 'AI_Legal_Portal_Presentation.pptx'")

if __name__ == "__main__":
    create_presentation()
